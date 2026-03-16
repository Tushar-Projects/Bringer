"""
Bringer RAG System — Document Loader

Responsible for ingesting various file formats from the documents directory
and converting them into clean, raw text.

Supported formats: PDF, DOCX, PPTX, TXT, MD.
"""

import os
from pathlib import Path
from typing import Optional, Dict, Any, List
from rich.console import Console

# Import parsers
import pypdf
import docx
import pptx
import markdown
from bs4 import BeautifulSoup

console = Console()

class DocumentLoader:
    def __init__(self):
        """Initialize the DocumentLoader."""
        # Mapping of extensions to their respective parsing methods
        self.parsers = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".txt": self._parse_txt,
            ".md": self._parse_md
        }

    def load_document(self, file_path: Path | str) -> Optional[List[Dict[str, Any]]]:
        """
        Loads a document and extracts its text along with system metadata (optionally per page).
        
        Args:
            file_path: Path to the document.
            
        Returns:
            List of dictionaries with 'content' and 'metadata', or None if parsing fails.
        """
        path = Path(file_path)
        
        if not path.exists():
            console.print(f"[red]Error: File not found - {path}[/red]")
            return None
            
        ext = path.suffix.lower()
        if ext not in self.parsers:
            console.print(f"[yellow]Warning: Unsupported file type - {ext}[/yellow]")
            return None
            
        try:
            pages = self.parsers[ext](path)
            if not pages:
                console.print(f"[yellow]Warning: No valid text extracted from {path.name}[/yellow]")
                return None
                
            return pages
            
        except Exception as e:
            console.print(f"[red]Error parsing {path.name}: {str(e)}[/red]")
            return None

    def _get_base_metadata(self, path: Path) -> Dict[str, Any]:
        """Helper to compute baseline metadata for a file."""
        return {
            "source_file": path.name,
            "file_path": str(path.as_posix()),
            "file_type": path.suffix.lower().lstrip('.')
        }

    def _parse_pdf(self, path: Path) -> List[Dict[str, Any]]:
        """Extract text from PDF files using pypdf, appending page numbers."""
        docs = []
        base_meta = self._get_base_metadata(path)
        
        with open(path, "rb") as file:
            reader = pypdf.PdfReader(file)
            for page_number, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    meta = base_meta.copy()
                    meta["page_number"] = page_number + 1
                    docs.append({
                        "content": page_text.strip(),
                        "metadata": meta
                    })
        return docs

    def _parse_docx(self, path: Path) -> List[Dict[str, Any]]:
        """Extract text from Word documents using python-docx."""
        doc = docx.Document(path)
        text = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        if text:
            return [{"content": "\n\n".join(text), "metadata": self._get_base_metadata(path)}]
        return []

    def _parse_pptx(self, path: Path) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint presentations using python-pptx (tracks slides as pages)."""
        prs = pptx.Presentation(path)
        docs = []
        base_meta = self._get_base_metadata(path)
        
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
            if text:
                meta = base_meta.copy()
                meta["page_number"] = i + 1
                docs.append({
                    "content": "\n\n".join(text),
                    "metadata": meta
                })
        return docs

    def _parse_txt(self, path: Path) -> List[Dict[str, Any]]:
        """Extract text from plain text files."""
        content = ""
        try:
            with open(path, "r", encoding="utf-8") as file:
                content = file.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as file:
                content = file.read()
                
        if content.strip():
            return [{"content": content.strip(), "metadata": self._get_base_metadata(path)}]
        return []

    def _parse_md(self, path: Path) -> List[Dict[str, Any]]:
        """Extract text from Markdown files, stripping HTML tags."""
        base = self._parse_txt(path)
        if not base:
            return []
            
        raw_text = base[0]["content"]
        html = markdown.markdown(raw_text)
        soup = BeautifulSoup(html, "html.parser")
        text = soup.get_text(separator="\n\n").strip()
        
        if text:
            base[0]["content"] = text
            return base
        return []


# Quick test trigger block (only runs if executed directly)
if __name__ == "__main__":
    import sys
    loader = DocumentLoader()
    if len(sys.argv) > 1:
        test_path = sys.argv[1]
        print(f"Testing loader on: {test_path}")
        pages = loader.load_document(test_path)
        if pages:
            print(f"\n--- Extracted Text Preview (First 500 chars of page 1) ---")
            print(pages[0]["content"][:500])
            print(f"\nTotal pages/sections extracted: {len(pages)}")
            print(f"Metadata of page 1: {pages[0]['metadata']}")
        else:
            print("Failed to extract text.")
    else:
        print("Usage: python document_loader.py <path_to_file>")
